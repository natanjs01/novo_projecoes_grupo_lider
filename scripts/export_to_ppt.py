#!/usr/bin/env python3
from __future__ import annotations

import argparse
import tempfile
from datetime import datetime
from pathlib import Path

from PIL import Image
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
SLIDES_DIR = ROOT / "site" / "public" / "slides"
OUTPUT_DIR = ROOT / "site" / "public"


def get_slide_files(first: int, last: int) -> list[Path]:
    files = []
    for path in SLIDES_DIR.glob("*.html"):
        prefix = path.name.split("_")[0]
        if prefix.isdigit() and first <= int(prefix) <= last:
            files.append(path)
    return sorted(files, key=lambda path: int(path.name.split("_")[0]))


def export_to_pptx(output_path: Path | None = None, first: int = 1, last: int = 22) -> str:
    slide_files = get_slide_files(first, last)
    if not slide_files:
        raise SystemExit("Nenhum slide encontrado no intervalo informado.")

    if output_path is None:
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        output_path = OUTPUT_DIR / f"Apresentacao_Grupo_Lider_{datetime.now():%Y%m%d_%H%M%S}.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory(prefix="slides_pptx_") as temp_dir:
        with sync_playwright() as playwright:
            browser = playwright.chromium.launch()
            page = browser.new_page(viewport={"width": 1920, "height": 1080}, device_scale_factor=2)
            page.emulate_media(media="screen")

            for index, slide_file in enumerate(slide_files, start=1):
                page.goto(slide_file.resolve().as_uri(), wait_until="networkidle")
                page.evaluate("document.fonts.ready")
                page.wait_for_timeout(1200)

                image_path = Path(temp_dir) / f"{index:02d}.png"
                page.screenshot(path=str(image_path), full_page=False)

                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    str(image_path),
                    left=0,
                    top=0,
                    width=prs.slide_width,
                    height=prs.slide_height,
                )
                print(f"✓ Slide {slide_file.name} exportado")

            browser.close()

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"\n✅ Apresentação exportada com sucesso: {output_path}")
    return str(output_path)


def main() -> None:
    parser = argparse.ArgumentParser(description="Exporta slides HTML para PPTX em 16:9.")
    parser.add_argument("--inicio", type=int, default=1, help="Número do primeiro slide")
    parser.add_argument("--fim", type=int, default=22, help="Número do último slide")
    parser.add_argument("--saida", type=Path, help="Caminho do PPTX de saída")
    args = parser.parse_args()

    export_to_pptx(args.saida, args.inicio, args.fim)


if __name__ == "__main__":
    main()